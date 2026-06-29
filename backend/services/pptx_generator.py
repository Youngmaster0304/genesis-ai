"""
Genesis AI — Professional PPTX Generator v2
Premium dark-themed investor pitch deck with charts, visual layouts, and rich styling.
"""
import io
import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ────────────────────────────────────────────────────
BG          = RGBColor(0x08, 0x0B, 0x14)
PANEL       = RGBColor(0x0F, 0x14, 0x21)
PANEL2      = RGBColor(0x16, 0x1D, 0x2F)
INDIGO      = RGBColor(0x63, 0x66, 0xF1)
PURPLE      = RGBColor(0x8B, 0x5C, 0xF6)
PINK        = RGBColor(0xEC, 0x48, 0x99)
EMERALD     = RGBColor(0x10, 0xB9, 0x81)
AMBER       = RGBColor(0xF5, 0x9E, 0x0B)
CYAN        = RGBColor(0x06, 0xB6, 0xD4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
SLATE200    = RGBColor(0xE2, 0xE8, 0xF0)
SLATE400    = RGBColor(0x94, 0xA3, 0xB8)
SLATE600    = RGBColor(0x47, 0x55, 0x69)
SLATE800    = RGBColor(0x1E, 0x29, 0x3B)

W  = Inches(13.33)
H  = Inches(7.5)
PHASE_COLORS = [INDIGO, PURPLE, PINK, AMBER]


# ─── Low-level primitives ─────────────────────────────────────

def _prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def _rect(slide, l, t, w, h, fill, radius=False):
    sp = slide.shapes.add_shape(1, l, t, w, h)
    sp.line.fill.background()
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    return sp

def _txt(slide, text, l, t, w, h, size=12, bold=False, italic=False,
         color=WHITE, align=PP_ALIGN.LEFT, wrap=True, font="Segoe UI"):
    if not text:
        return None
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = wrap
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name   = font
    return box

def _tag(slide, text, l, t, color=INDIGO, size=7.5):
    """Small pill label."""
    bw = Inches(len(text) * 0.085 + 0.25)
    _rect(slide, l, t, bw, Inches(0.22), _darken(color, 0.25))
    _txt(slide, text.upper(), l + Inches(0.08), t + Pt(1), bw, Inches(0.22),
         size=size, bold=True, color=color)
    return bw

def _darken(color: RGBColor, factor=0.4):
    return RGBColor(
        int(color[0] * factor),
        int(color[1] * factor),
        int(color[2] * factor),
    )

def _bar_top(slide, color=INDIGO, h=Pt(4)):
    _rect(slide, 0, 0, W, h, color)

def _gradient_stripe(slide):
    """Fake gradient using overlapping semi-transparent strips."""
    for i, col in enumerate([INDIGO, PURPLE, PINK]):
        sw = Inches(13.33 / 3)
        _rect(slide, i * sw, 0, sw, Pt(4), col)

def _slide_num(slide, n, total=12):
    _txt(slide, f"{n:02}/{total}",
         Inches(12.2), Inches(7.12), Inches(1), Inches(0.3),
         size=7, color=SLATE600, align=PP_ALIGN.RIGHT)

def _logo(slide):
    _txt(slide, "GENESIS AI",
         Inches(0.35), Inches(7.12), Inches(2), Inches(0.3),
         size=7, bold=True, color=SLATE600)

def _divider(slide, top, color=SLATE800):
    _rect(slide, Inches(0.5), top, Inches(12.33), Pt(1), color)

def _section_header(slide, title, subtitle="", icon="", top=Inches(0.38)):
    """Consistent section title row."""
    if icon:
        _txt(slide, icon, Inches(0.5), top, Inches(0.5), Inches(0.55), size=20)
        tx = Inches(1.1)
    else:
        tx = Inches(0.5)
    _txt(slide, title, tx, top, Inches(10), Inches(0.48),
         size=24, bold=True, color=WHITE, font="Segoe UI Semibold")
    if subtitle:
        _txt(slide, subtitle, tx, top + Inches(0.48), Inches(10), Inches(0.3),
             size=9, color=SLATE400)

def _score_gauge(slide, label, value, x, y, w=Inches(1.9)):
    """Mini score gauge: filled bar + number."""
    pct  = min(max(int(value), 0), 100) / 100
    bar_h = Inches(0.06)
    bar_w = w - Inches(0.1)
    _txt(slide, label, x, y, w, Inches(0.22), size=7.5, color=SLATE400)
    _txt(slide, str(value), x, y + Inches(0.22), w, Inches(0.4),
         size=22, bold=True, color=INDIGO)
    _rect(slide, x, y + Inches(0.65), bar_w, bar_h, SLATE800)
    _rect(slide, x, y + Inches(0.65), Inches((w.inches - 0.1) * pct), bar_h, INDIGO)


# ─── Slide builders ───────────────────────────────────────────

def _s1_title(prs, report):
    """Slide 1 — Title + Hook"""
    slide = _blank(prs)
    _bg(slide)

    # Left dark panel
    _rect(slide, 0, 0, Inches(8.5), H, PANEL)
    # Gradient strip
    _gradient_stripe(slide)

    # Big title
    title = (report.get("input_summary") or "Innovation R&D Report")[:90]
    _txt(slide, title,
         Inches(0.6), Inches(1.3), Inches(7.6), Inches(3.0),
         size=36, bold=True, color=WHITE, font="Segoe UI Semibold", wrap=True)

    # Hook
    hook = (report.get("sections", {})
                  .get("investor_pitch", {})
                  .get("hook", "Transforming an industry with AI-powered innovation."))
    _txt(slide, f'"{hook}"',
         Inches(0.6), Inches(4.5), Inches(7.6), Inches(1.2),
         size=14, italic=True, color=INDIGO)

    # Meta
    date = (report.get("created_at", "")[:10] or
            datetime.datetime.utcnow().strftime("%Y-%m-%d"))
    _txt(slide, f"Genesis AI  ·  {date}",
         Inches(0.6), Inches(5.9), Inches(7), Inches(0.3),
         size=8.5, color=SLATE400)

    # Right score panel
    scores = report.get("scores", {})
    _txt(slide, "INNOVATION SCORES",
         Inches(9.0), Inches(1.2), Inches(3.9), Inches(0.3),
         size=8, bold=True, color=SLATE400)
    _divider(slide, Inches(1.6), SLATE800)

    gauge_data = [
        ("Overall",    scores.get("overall",              78)),
        ("Innovation", scores.get("innovation",           80)),
        ("Market",     scores.get("market_opportunity",   76)),
        ("Technical",  scores.get("technical_feasibility",82)),
        ("Business",   scores.get("business_viability",   74)),
        ("IP Novelty", scores.get("patent_novelty",       77)),
    ]
    gy = Inches(1.75)
    for label, val in gauge_data:
        _score_gauge(slide, label, val, Inches(9.0), gy)
        gy += Inches(0.88)

    _slide_num(slide, 1)
    _logo(slide)


def _s2_problem(prs, report):
    slide = _blank(prs)
    _bg(slide)
    _rect(slide, 0, 0, W, Pt(4), PINK)
    _section_header(slide, "The Problem", "Why this pain point demands a new solution", "🔴")

    pitch = report.get("sections", {}).get("investor_pitch", {})
    prob  = pitch.get("problem_statement", "The current market solutions are fragmented, expensive, and unable to serve modern enterprise needs at scale.")

    # Big problem statement
    _rect(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.6), PANEL)
    _rect(slide, Inches(0.5), Inches(1.3), Pt(5), Inches(1.6), PINK)
    _txt(slide, prob, Inches(0.75), Inches(1.45), Inches(11.8), Inches(1.4),
         size=14, color=SLATE200, italic=True)

    # Pain points from executive summary
    exec_s = report.get("sections", {}).get("executive_summary", "")
    excerpt = exec_s[:450] + "…" if len(exec_s) > 450 else exec_s
    _txt(slide, "MARKET PAIN EVIDENCE", Inches(0.5), Inches(3.2),
         Inches(4), Inches(0.25), size=7.5, bold=True, color=SLATE400)
    _rect(slide, Inches(0.5), Inches(3.5), Inches(12.3), Inches(3.35), PANEL2)
    _txt(slide, excerpt, Inches(0.75), Inches(3.65), Inches(11.8), Inches(3.1),
         size=10.5, color=SLATE200)

    _slide_num(slide, 2)
    _logo(slide)


def _s3_solution(prs, report):
    slide = _blank(prs)
    _bg(slide)
    _rect(slide, 0, 0, W, Pt(4), EMERALD)
    _section_header(slide, "The Solution", "Our AI-powered approach and core differentiators", "✅")

    pitch = report.get("sections", {}).get("investor_pitch", {})
    sol   = pitch.get("solution_statement", "A multi-agent AI platform that automates end-to-end R&D innovation pipeline.")
    _rect(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.4), PANEL)
    _rect(slide, Inches(0.5), Inches(1.3), Pt(5), Inches(1.4), EMERALD)
    _txt(slide, sol, Inches(0.75), Inches(1.45), Inches(11.8), Inches(1.2),
         size=14, color=SLATE200, italic=True)

    # Differentiators
    gaps = (report.get("sections", {})
                  .get("novelty_score_breakdown", {})
                  .get("gaps", []))
    _txt(slide, "KEY DIFFERENTIATORS", Inches(0.5), Inches(3.0),
         Inches(6), Inches(0.25), size=7.5, bold=True, color=SLATE400)

    icons = ["✦", "◆", "●", "★"]
    diff_y = Inches(3.35)
    for i, gap in enumerate(gaps[:4]):
        g = str(gap)[:120]
        _rect(slide, Inches(0.5), diff_y, Inches(12.3), Inches(0.65), PANEL if i % 2 == 0 else PANEL2)
        _txt(slide, icons[i % len(icons)], Inches(0.65), diff_y + Pt(4),
             Inches(0.3), Inches(0.6), size=14, color=EMERALD)
        _txt(slide, g, Inches(1.05), diff_y + Pt(5),
             Inches(11.5), Inches(0.55), size=11, color=SLATE200)
        diff_y += Inches(0.72)

    _slide_num(slide, 3)
    _logo(slide)


def _s4_market(prs, report):
    slide = _blank(prs)
    _bg(slide)
    _rect(slide, 0, 0, W, Pt(4), AMBER)
    _section_header(slide, "Market Opportunity", "TAM · SAM · SOM with macro trend analysis", "📈")

    market = report.get("sections", {}).get("market_analysis", {})
    ms     = market.get("market_size", "TAM: $42B — SAM: $8B — SOM: $1.2B (CAGR 28%)")
    _txt(slide, ms, Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.7),
         size=18, bold=True, color=AMBER)

    # Visual TAM/SAM/SOM bars
    bars = [("TAM", 1.0, INDIGO), ("SAM", 0.55, PURPLE), ("SOM", 0.18, PINK)]
    bar_top = Inches(2.2)
    bar_max_w = Inches(10)
    for label, pct, col in bars:
        _rect(slide, Inches(0.5), bar_top, bar_max_w, Inches(0.38), SLATE800)
        _rect(slide, Inches(0.5), bar_top, Inches(10 * pct), Inches(0.38), col)
        _txt(slide, label, Inches(10.7), bar_top, Inches(1.2), Inches(0.38),
             size=10, bold=True, color=col)
        bar_top += Inches(0.54)

    # Trends
    trends = market.get("trends", [])
    _txt(slide, "MACRO TRENDS", Inches(0.5), Inches(3.75),
         Inches(6), Inches(0.25), size=7.5, bold=True, color=SLATE400)
    trend_y = Inches(4.05)
    for i, t in enumerate(trends[:4]):
        txt = str(t)[:130]
        _rect(slide, Inches(0.5), trend_y, Inches(12.3), Inches(0.6), PANEL if i % 2 == 0 else PANEL2)
        _txt(slide, f"→  {txt}", Inches(0.7), trend_y + Pt(4),
             Inches(12.0), Inches(0.5), size=10, color=SLATE200)
        trend_y += Inches(0.65)

    # Why now
    timing = market.get("market_timing", "")
    if timing and trend_y < Inches(7.1):
        _rect(slide, Inches(0.5), trend_y, Inches(12.3), Inches(0.65), _darken(AMBER, 0.3))
        _txt(slide, f"⏱  WHY NOW: {timing[:160]}",
             Inches(0.7), trend_y + Pt(4), Inches(12.0), Inches(0.55),
             size=10, italic=True, color=AMBER)

    _slide_num(slide, 4)
    _logo(slide)


def _s5_competition(prs, report):
    slide = _blank(prs)
    _bg(slide)
    _rect(slide, 0, 0, W, Pt(4), PURPLE)
    _section_header(slide, "Competitive Landscape", "Named competitors with exploitable weaknesses", "🛡️")

    comps = report.get("sections", {}).get("competitor_analysis", [])
    col_x = [Inches(0.4), Inches(2.5), Inches(6.4), Inches(10.1)]
    col_w = [Inches(1.9), Inches(3.7), Inches(3.5), Inches(3.0)]
    headers = ["Company", "Strength", "Weakness", "Type / Funding"]

    # Header row
    _rect(slide, Inches(0.4), Inches(1.35), Inches(12.5), Inches(0.4), INDIGO)
    for i, h in enumerate(headers):
        _txt(slide, h, col_x[i] + Inches(0.08), Inches(1.38), col_w[i], Inches(0.38),
             size=9, bold=True, color=WHITE)

    row_y = Inches(1.78)
    for idx, c in enumerate(comps[:5]):
        bg = PANEL if idx % 2 == 0 else PANEL2
        _rect(slide, Inches(0.4), row_y, Inches(12.5), Inches(0.88), bg)
        # Left accent
        type_col = PINK if c.get("type") == "Direct" else (AMBER if c.get("type") == "Emerging" else INDIGO)
        _rect(slide, Inches(0.4), row_y, Pt(4), Inches(0.88), type_col)

        name = str(c.get("name", f"Comp {idx+1}"))[:28]
        strength = str(c.get("strength", ""))[:80]
        weakness = str(c.get("weakness", ""))[:80]
        funding  = f"{c.get('type','?')} · {c.get('funding','')}"[:30]

        _txt(slide, name, col_x[0] + Inches(0.1), row_y + Pt(4), col_w[0], Inches(0.82),
             size=9.5, bold=True, color=INDIGO)
        _txt(slide, strength, col_x[1] + Inches(0.08), row_y + Pt(4), col_w[1], Inches(0.82),
             size=8.5, color=SLATE200)
        _txt(slide, weakness, col_x[2] + Inches(0.08), row_y + Pt(4), col_w[2], Inches(0.82),
             size=8.5, color=PINK)
        _txt(slide, funding, col_x[3] + Inches(0.08), row_y + Pt(4), col_w[3], Inches(0.82),
             size=8, color=SLATE400)
        row_y += Inches(0.92)

    # Our edge
    edge_y = row_y + Inches(0.08)
    if edge_y < Inches(7.1):
        _rect(slide, Inches(0.4), edge_y, Inches(12.5), Inches(0.7), _darken(EMERALD, 0.25))
        breakthrough = (report.get("sections", {})
                              .get("novelty_score_breakdown", {})
                              .get("breakthrough_classification", "Disruptive innovation."))
        _txt(slide, f"⚡  OUR EDGE:  {str(breakthrough)[:180]}",
             Inches(0.65), edge_y + Pt(5), Inches(12.0), Inches(0.6),
             size=10, bold=False, color=EMERALD)

    _slide_num(slide, 5)
    _logo(slide)


def _s6_technology(prs, report):
    slide = _blank(prs)
    _bg(slide)
    _rect(slide, 0, 0, W, Pt(4), CYAN)
    _section_header(slide, "Technology & Architecture", "System design, stack, and scalability blueprint", "🏗️")

    arch = report.get("sections", {}).get("technical_architecture", {})
    summ = str(arch.get("architecture_summary", "A cloud-native multi-agent AI system."))[:400]
    _rect(slide, Inches(0.5), Inches(1.35), Inches(12.3), Inches(1.1), PANEL)
    _rect(slide, Inches(0.5), Inches(1.35), Pt(5), Inches(1.1), CYAN)
    _txt(slide, summ, Inches(0.75), Inches(1.45), Inches(11.8), Inches(1.0),
         size=10.5, color=SLATE200)

    # Architecture nodes as flow boxes
    nodes = arch.get("diagram_nodes", [])
    if nodes:
        _txt(slide, "SYSTEM COMPONENTS", Inches(0.5), Inches(2.65),
             Inches(5), Inches(0.25), size=7.5, bold=True, color=SLATE400)
        cols = 4
        bw = Inches(2.9)
        bh = Inches(0.55)
        gx = Inches(0.25)
        start_x = Inches(0.5)
        node_y  = Inches(2.95)
        for i, nd in enumerate(nodes[:8]):
            col = i % cols
            row = i // cols
            nx = start_x + col * (bw + gx)
            ny = node_y  + row * (bh + Inches(0.12))
            nc = [INDIGO, PURPLE, PINK, AMBER, EMERALD, CYAN, INDIGO, PURPLE][i % 8]
            _rect(slide, nx, ny, bw, bh, PANEL2)
            _rect(slide, nx, ny, Pt(4), bh, nc)
            _txt(slide, str(nd)[:30], nx + Inches(0.12), ny + Pt(4),
                 bw - Inches(0.18), bh, size=9.5, color=SLATE200)

    # Tech stack badges row
    stack = report.get("sections", {}).get("technology_stack", [])
    if stack:
        sy = Inches(5.85)
        _txt(slide, "CORE STACK", Inches(0.5), sy, Inches(3), Inches(0.25),
             size=7.5, bold=True, color=SLATE400)
        sx = Inches(0.5)
        sy += Inches(0.3)
        for tech in stack[:12]:
            tw = Pt(len(str(tech.get("name", ""))) * 6 + 24)
            if sx + tw > Inches(12.8):
                sx  = Inches(0.5)
                sy += Inches(0.4)
            cat = str(tech.get("category", "Core"))
            col_map = {
                "Frontend": INDIGO, "Backend": PURPLE, "ML": PINK,
                "Infra": AMBER, "DB": EMERALD, "Infrastructure": AMBER,
                "Database": EMERALD, "Machine Learning": PINK,
            }
            tc = col_map.get(cat, INDIGO)
            _rect(slide, sx, sy, tw, Inches(0.3), _darken(tc, 0.3))
            _txt(slide, str(tech.get("name", ""))[:18], sx + Pt(6), sy + Pt(2),
                 tw, Inches(0.28), size=8, bold=True, color=tc)
            sx += tw + Inches(0.12)

    _slide_num(slide, 6)
    _logo(slide)


def _s7_business(prs, report):
    slide = _blank(prs)
    _bg(slide)
    _rect(slide, 0, 0, W, Pt(4), EMERALD)
    _section_header(slide, "Business Model", "Revenue architecture and unit economics", "💼")

    biz = report.get("sections", {}).get("business_model", {})
    desc = str(biz.get("description", "B2B SaaS with tiered pricing."))[:320]
    _rect(slide, Inches(0.5), Inches(1.35), Inches(12.3), Inches(1.0), PANEL)
    _rect(slide, Inches(0.5), Inches(1.35), Pt(5), Inches(1.0), EMERALD)
    _txt(slide, desc, Inches(0.75), Inches(1.45), Inches(11.8), Inches(0.9),
         size=10.5, color=SLATE200)

    # Unit economics bar
    ue = str(biz.get("unit_economics", ""))
    if ue:
        _rect(slide, Inches(0.5), Inches(2.55), Inches(12.3), Inches(0.5), _darken(EMERALD, 0.28))
        _txt(slide, f"📊  {ue[:160]}",
             Inches(0.7), Inches(2.62), Inches(12.0), Inches(0.42),
             size=10, color=EMERALD, bold=True)

    # Revenue projections
    rev = str(biz.get("revenue_projections", "Year 1: $1.2M  ·  Year 2: $4.5M  ·  Year 3: $14M"))
    _txt(slide, "REVENUE PROJECTIONS", Inches(0.5), Inches(3.25),
         Inches(6), Inches(0.25), size=7.5, bold=True, color=SLATE400)
    _txt(slide, rev, Inches(0.5), Inches(3.55), Inches(12.3), Inches(0.55),
         size=14, bold=True, color=AMBER)

    # Pricing tier cards
    tiers = biz.get("tiers", [])
    _txt(slide, "PRICING TIERS", Inches(0.5), Inches(4.25),
         Inches(6), Inches(0.25), size=7.5, bold=True, color=SLATE400)
    tier_cols = [INDIGO, PURPLE, PINK]
    tw = Inches(3.95)
    tx = Inches(0.5)
    for i, tier in enumerate(tiers[:3]):
        tc = tier_cols[i % len(tier_cols)]
        _rect(slide, tx, Inches(4.55), tw, Inches(2.5), PANEL2)
        _rect(slide, tx, Inches(4.55), tw, Pt(4), tc)
        _txt(slide, f"TIER {i+1}", tx + Inches(0.15), Inches(4.62),
             tw - Inches(0.2), Inches(0.28), size=7.5, bold=True, color=tc)
        _txt(slide, str(tier)[:140], tx + Inches(0.15), Inches(4.95),
             tw - Inches(0.25), Inches(2.0), size=9.5, color=SLATE200)
        tx += tw + Inches(0.22)

    _slide_num(slide, 7)
    _logo(slide)


def _s8_financials(prs, report):
    slide = _blank(prs)
    _bg(slide)
    _rect(slide, 0, 0, W, Pt(4), AMBER)
    _section_header(slide, "Financial Projections", "3-year ARR forecast with key assumptions", "📊")

    biz = report.get("sections", {}).get("business_model", {})
    rev = str(biz.get("revenue_projections", "Year 1: $1.2M ARR  ·  Year 2: $4.5M ARR  ·  Year 3: $14M ARR"))
    _txt(slide, rev, Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.6),
         size=16, bold=True, color=AMBER)

    # Bar chart
    chart_bottom = Inches(6.9)
    chart_left   = Inches(1.5)
    bar_data = [
        ("Year 1", 0.22, INDIGO),
        ("Year 2", 0.55, PURPLE),
        ("Year 3", 1.00, PINK),
    ]
    max_h = Inches(4.6)
    bw    = Inches(2.4)
    gap   = Inches(1.0)

    # Y-axis labels
    for pct, lbl in [(0, "0"), (0.5, "50%"), (1.0, "100%")]:
        y = chart_bottom - max_h * pct
        _txt(slide, lbl, Inches(0.4), y - Inches(0.15), Inches(0.9), Inches(0.3),
             size=8, color=SLATE400, align=PP_ALIGN.RIGHT)
        _rect(slide, chart_left, y, Inches(9.5), Pt(1), SLATE800)

    for i, (yr, pct, col) in enumerate(bar_data):
        bx = chart_left + i * (bw + gap)
        bh = max_h * pct
        by = chart_bottom - bh
        # Shadow
        _rect(slide, bx + Inches(0.08), by + Inches(0.08), bw, bh, _darken(col, 0.35))
        # Bar
        _rect(slide, bx, by, bw, bh, col)
        # Value label on bar
        _txt(slide, yr, bx, chart_bottom + Pt(6), bw, Inches(0.3),
             size=11, bold=True, color=SLATE400, align=PP_ALIGN.CENTER)

    # Baseline
    _rect(slide, chart_left, chart_bottom, Inches(9.5), Pt(2), SLATE600)

    _slide_num(slide, 8)
    _logo(slide)


def _s9_gtm(prs, report):
    slide = _blank(prs)
    _bg(slide)
    _rect(slide, 0, 0, W, Pt(4), PURPLE)
    _section_header(slide, "Go-To-Market Strategy", "Beachhead, acquisition channels, and first 100 customers", "🚀")

    gtm = report.get("sections", {}).get("gtm_strategy", {})

    # Beachhead
    beachhead = str(gtm.get("beachhead", "Focus on early adopter enterprise segment."))[:220]
    _rect(slide, Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.75), _darken(INDIGO, 0.35))
    _txt(slide, f"🎯  BEACHHEAD:  {beachhead}",
         Inches(0.7), Inches(1.43), Inches(12.0), Inches(0.65), size=11, color=SLATE200)

    # Channels as funnel cards
    channels = gtm.get("channels", [])
    _txt(slide, "ACQUISITION FUNNEL", Inches(0.5), Inches(2.3),
         Inches(6), Inches(0.25), size=7.5, bold=True, color=SLATE400)
    ch_cols = [INDIGO, PURPLE, PINK, AMBER]
    ch_y = Inches(2.62)
    ch_widths = [Inches(12.3), Inches(10.5), Inches(8.8), Inches(7.2)]
    for i, ch in enumerate(channels[:4]):
        cw = ch_widths[i]
        cx = Inches(0.5) + (Inches(12.3) - cw) / 2
        col = ch_cols[i % len(ch_cols)]
        _rect(slide, cx, ch_y, cw, Inches(0.6), _darken(col, 0.3))
        _rect(slide, cx, ch_y, Pt(4), Inches(0.6), col)
        _txt(slide, f"  {str(ch)[:120]}", cx + Inches(0.12), ch_y + Pt(6),
             cw - Inches(0.2), Inches(0.52), size=10, color=SLATE200)
        ch_y += Inches(0.68)

    # First 100 customers
    f100 = str(gtm.get("first_100_customers", ""))[:220]
    if f100 and ch_y < Inches(7.0):
        _rect(slide, Inches(0.5), ch_y + Inches(0.1), Inches(12.3), Inches(0.72), PANEL2)
        _txt(slide, f"👥  FIRST 100:  {f100}",
             Inches(0.7), ch_y + Inches(0.17), Inches(12.0), Inches(0.6),
             size=10, color=SLATE200)

    _slide_num(slide, 9)
    _logo(slide)


def _s10_ip(prs, report):
    slide = _blank(prs)
    _bg(slide)
    _rect(slide, 0, 0, W, Pt(4), EMERALD)
    _section_header(slide, "IP & Patent Landscape", "Freedom-to-operate analysis and IP moat", "📜")

    patent = report.get("sections", {}).get("patent_landscape", {})
    moat   = str(patent.get("ip_moat_rating", "Moderate")).split()[0]
    moat_col = EMERALD if "Strong" in moat else (AMBER if "Moderate" in moat else PINK)

    # Moat badge
    _rect(slide, Inches(0.5), Inches(1.35), Inches(2.5), Inches(0.7), _darken(moat_col, 0.3))
    _txt(slide, f"IP MOAT: {moat}",
         Inches(0.6), Inches(1.45), Inches(2.4), Inches(0.55),
         size=13, bold=True, color=moat_col, align=PP_ALIGN.CENTER)

    status = str(patent.get("status", "Freedom-to-operate confirmed."))[:300]
    _txt(slide, status, Inches(3.2), Inches(1.42), Inches(9.8), Inches(0.9),
         size=10.5, color=SLATE200)

    # Novelty gaps
    gaps = patent.get("novelty_gaps", [])
    _txt(slide, "CLAIMABLE WHITE-SPACE", Inches(0.5), Inches(2.55),
         Inches(6), Inches(0.25), size=7.5, bold=True, color=SLATE400)
    gy = Inches(2.88)
    for i, g in enumerate(gaps[:4]):
        _rect(slide, Inches(0.5), gy, Inches(12.3), Inches(0.58),
              PANEL if i % 2 == 0 else PANEL2)
        _txt(slide, f"✓  {str(g)[:140]}", Inches(0.7), gy + Pt(5),
             Inches(12.0), Inches(0.5), size=10.5, color=EMERALD)
        gy += Inches(0.64)

    # Claim strategy
    strat = str(patent.get("claim_strategy", ""))[:250]
    if strat and gy < Inches(7.0):
        _rect(slide, Inches(0.5), gy + Inches(0.1), Inches(12.3), Inches(0.75), _darken(EMERALD, 0.25))
        _txt(slide, f"📋  FILING STRATEGY:  {strat}",
             Inches(0.7), gy + Inches(0.18), Inches(12.0), Inches(0.65),
             size=10, color=EMERALD)

    _slide_num(slide, 10)
    _logo(slide)


def _s11_roadmap(prs, report):
    slide = _blank(prs)
    _bg(slide)
    _rect(slide, 0, 0, W, Pt(4), INDIGO)
    _section_header(slide, "MVP Roadmap", "Phase-by-phase build plan with success metrics", "🗺️")

    phases = report.get("sections", {}).get("mvp_roadmap", [])
    if not phases:
        _txt(slide, "Roadmap data not available for this report.",
             Inches(0.5), Inches(2.5), Inches(12), Inches(1),
             size=14, color=SLATE400)
        _slide_num(slide, 11)
        _logo(slide)
        return

    n = min(len(phases), 4)
    pw = (Inches(12.3) - Inches(0.18) * (n - 1)) / n
    px = Inches(0.5)
    p_colors = [INDIGO, PURPLE, PINK, AMBER]

    for i, ph in enumerate(phases[:n]):
        col = p_colors[i % len(p_colors)]
        # Phase header
        _rect(slide, px, Inches(1.35), pw, Inches(0.55), col)
        ph_name = str(ph.get("phase", f"Phase {i+1}"))[:30]
        _txt(slide, ph_name, px + Inches(0.1), Inches(1.4),
             pw - Inches(0.15), Inches(0.48), size=9.5, bold=True, color=WHITE)
        # Duration badge
        dur = str(ph.get("duration", ""))[:20]
        _rect(slide, px, Inches(1.92), pw, Inches(0.32), _darken(col, 0.45))
        _txt(slide, dur, px + Inches(0.1), Inches(1.94),
             pw - Inches(0.15), Inches(0.28), size=8.5, bold=True, color=col)
        # Body
        _rect(slide, px, Inches(2.27), pw, Inches(4.55), PANEL2)
        # Goal
        goal = str(ph.get("goal", ""))[:70]
        if goal:
            _txt(slide, f"🎯 {goal}", px + Inches(0.1), Inches(2.35),
                 pw - Inches(0.15), Inches(0.55), size=8.5, italic=True, color=SLATE400)
        # Tasks
        ty = Inches(2.95)
        for task in (ph.get("tasks") or [])[:5]:
            t = str(task)[:60]
            _txt(slide, f"◆  {t}", px + Inches(0.1), ty,
                 pw - Inches(0.15), Inches(0.5), size=8.5, color=SLATE200)
            ty += Inches(0.52)
        # Success metric
        sm = str(ph.get("success_metric", ""))[:70]
        if sm:
            _rect(slide, px, Inches(6.6), pw, Inches(0.55), _darken(col, 0.32))
            _txt(slide, f"✓ {sm}", px + Inches(0.08), Inches(6.65),
                 pw - Inches(0.14), Inches(0.48), size=8, color=col)

        px += pw + Inches(0.18)

    _slide_num(slide, 11)
    _logo(slide)


def _s12_ask(prs, report):
    slide = _blank(prs)
    _bg(slide)
    _gradient_stripe(slide)
    _section_header(slide, "The Ask", "Investment opportunity and use of funds", "💰")

    pitch = report.get("sections", {}).get("investor_pitch", {})
    hook  = str(pitch.get("hook", ""))[:200]
    if hook:
        _rect(slide, Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.75), _darken(INDIGO, 0.35))
        _rect(slide, Inches(0.5), Inches(1.35), Pt(5), Inches(0.75), INDIGO)
        _txt(slide, f'"{hook}"',
             Inches(0.75), Inches(1.43), Inches(12.0), Inches(0.65),
             size=13, italic=True, color=INDIGO)

    traction = str(pitch.get("traction_needed", ""))[:300]
    if traction:
        _rect(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.0), PANEL2)
        _txt(slide, traction, Inches(0.7), Inches(2.4), Inches(12.0), Inches(0.88),
             size=11, color=SLATE200)

    # Score summary grid
    scores = report.get("scores", {})
    _txt(slide, "GENESIS AI VALIDATION SCORES",
         Inches(0.5), Inches(3.55), Inches(12), Inches(0.28),
         size=7.5, bold=True, color=SLATE400)

    gauge_data = [
        ("Overall",    scores.get("overall", 0),              INDIGO),
        ("Innovation", scores.get("innovation", 0),           PURPLE),
        ("Market",     scores.get("market_opportunity", 0),   AMBER),
        ("Technical",  scores.get("technical_feasibility", 0),CYAN),
        ("Business",   scores.get("business_viability", 0),   EMERALD),
        ("IP Novelty", scores.get("patent_novelty", 0),       PINK),
    ]
    sw = Inches(1.9)
    sx = Inches(0.5)
    sy = Inches(3.9)
    for label, val, col in gauge_data:
        _rect(slide, sx, sy, sw, Inches(1.55), PANEL)
        _txt(slide, label, sx + Inches(0.1), sy + Inches(0.1),
             sw - Inches(0.15), Inches(0.28), size=7.5, color=SLATE400)
        _txt(slide, str(val), sx + Inches(0.1), sy + Inches(0.38),
             sw - Inches(0.15), Inches(0.7), size=30, bold=True, color=col)
        # Mini bar
        pct = min(max(int(val), 0), 100) / 100
        _rect(slide, sx + Inches(0.1), sy + Inches(1.15), sw - Inches(0.2), Inches(0.08), SLATE800)
        _rect(slide, sx + Inches(0.1), sy + Inches(1.15), Inches((sw.inches - 0.2) * pct), Inches(0.08), col)
        sx += sw + Inches(0.25)

    # Footer
    _txt(slide, "genesis-ai.io  ·  Powered by Gemma 4 31B via Cerebras  ·  Confidential & Proprietary",
         Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.3),
         size=7, color=SLATE600, align=PP_ALIGN.CENTER)

    _slide_num(slide, 12, 12)
    _logo(slide)


# ─── Public API ───────────────────────────────────────────────

def generate_pptx(report: dict) -> bytes:
    """
    Generate a professional 12-slide investor pitch deck from Genesis AI report.
    Returns raw .pptx bytes ready for HTTP response.
    """
    prs = _prs()
    _s1_title(prs, report)
    _s2_problem(prs, report)
    _s3_solution(prs, report)
    _s4_market(prs, report)
    _s5_competition(prs, report)
    _s6_technology(prs, report)
    _s7_business(prs, report)
    _s8_financials(prs, report)
    _s9_gtm(prs, report)
    _s10_ip(prs, report)
    _s11_roadmap(prs, report)
    _s12_ask(prs, report)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
